"""
Add slide 3 to curriculum_slide.pptx
Reproduces the image from slide 2 as actual editable text and shapes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy
from lxml import etree

# ── helpers ──────────────────────────────────────────────────────────────────

def rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def add_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox, tf

def set_cell_bg(shape, color_rgb):
    """Set solid fill on a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def set_no_fill(shape):
    shape.fill.background()

def set_line(shape, color_rgb=None, width_pt=1.5):
    ln = shape.line
    if color_rgb:
        ln.color.rgb = color_rgb
        ln.width = Pt(width_pt)
    else:
        ln.fill.background()  # no line

def add_paragraph(tf, text='', bold=False, size_pt=10,
                  color_hex='#FFFFFF', align=PP_ALIGN.LEFT,
                  space_before_pt=0, font_name='Trebuchet MS',
                  level=0, italic=False):
    """Add a single-run paragraph."""
    p = tf.add_paragraph()
    p.space_before = Pt(space_before_pt)
    p.alignment = align
    p.level = level
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(size_pt)
    run.font.color.rgb = rgb(color_hex)
    run.font.name = font_name
    return p

def add_mixed_paragraph(tf, bold_text, bold_color, bold_size,
                         normal_text, normal_color, normal_size,
                         space_before_pt=4, align=PP_ALIGN.LEFT,
                         font_name='Trebuchet MS'):
    """Bold run followed by normal run."""
    p = tf.add_paragraph()
    p.space_before = Pt(space_before_pt)
    p.alignment = align
    if bold_text:
        r1 = p.add_run()
        r1.text = bold_text
        r1.font.bold = True
        r1.font.size = Pt(bold_size)
        r1.font.color.rgb = rgb(bold_color)
        r1.font.name = font_name
    if normal_text:
        r2 = p.add_run()
        r2.text = normal_text
        r2.font.bold = False
        r2.font.size = Pt(normal_size)
        r2.font.color.rgb = rgb(normal_color)
        r2.font.name = font_name
    return p

def add_continuation(tf, text, color_hex, size_pt=5.5, font_name='Trebuchet MS'):
    p = tf.add_paragraph()
    p.space_before = Pt(0)
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = '  ' + text
    run.font.bold = False
    run.font.size = Pt(size_pt)
    run.font.color.rgb = rgb(color_hex)
    run.font.name = font_name
    return p

def rounded_rect(slide, left, top, width, height,
                 fill_color=None, line_color=None, line_pt=1.5,
                 corner_radius_emu=None):
    """Add a rounded rectangle (MSO_SHAPE_TYPE 5 = ROUNDED_RECTANGLE)."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        left, top, width, height
    )
    # corner radius adjustment
    if corner_radius_emu is not None:
        # prstGeom -> avLst -> gd name="adj"
        sp_tree = shape._element
        prstGeom = sp_tree.find('.//' + qn('a:prstGeom'))
        if prstGeom is not None:
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            # remove existing
            for gd in avLst.findall(qn('a:gd')):
                avLst.remove(gd)
            gd = etree.SubElement(avLst, qn('a:gd'))
            # adj value: ratio of min(w,h)/2; EMU-based pct * 100000
            # Use a small fixed value like 20000 (~2% of shape size)
            gd.set('name', 'adj')
            gd.set('fmla', f'val {corner_radius_emu}')
    if fill_color:
        set_cell_bg(shape, fill_color)
    else:
        set_no_fill(shape)
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


# ── Main ─────────────────────────────────────────────────────────────────────

prs = Presentation(r"C:/Users/hiond/workspace/dmap/curriculum_slide.pptx")

# Verify slide dimensions
print(f"Slide width:  {prs.slide_width.inches:.3f} in")
print(f"Slide height: {prs.slide_height.inches:.3f} in")
print(f"Current slide count: {len(prs.slides)}")

slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

# ── Background ────────────────────────────────────────────────────────────────
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = rgb('#2D2D2D')

# ── Title ─────────────────────────────────────────────────────────────────────
W = Inches(10)
title_left = Inches(0.15)
title_top  = Inches(0.08)
title_w    = Inches(9.70)
title_h    = Inches(0.55)

txBox, tf = add_textbox(slide, title_left, title_top, title_w, title_h)
tf.auto_size = None
# remove default empty paragraph
for _ in range(len(tf.paragraphs)):
    tf._txBody.remove(tf._txBody.findall(qn('a:p'))[-1])

# Line 1
p1 = tf.add_paragraph()
p1.alignment = PP_ALIGN.CENTER
p1.space_before = Pt(0)
r1 = p1.add_run()
r1.text = "YET AI Agent 부트캠프 — 5일 커리큘럼 연관관계도"
r1.font.bold = True
r1.font.size = Pt(18)
r1.font.color.rgb = rgb('#FFFFFF')
r1.font.name = 'Trebuchet MS'

# Line 2
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(2)
r2 = p2.add_run()
r2.text = "대화  ▶  데이터  ▶  도구  ▶  지식  ▶  에이전트"
r2.font.bold = False
r2.font.size = Pt(10)
r2.font.color.rgb = rgb('#22D3EE')
r2.font.name = 'Trebuchet MS'

# ── Column layout constants ────────────────────────────────────────────────────
COL_LEFT_MARGIN = Inches(0.20)
COL_WIDTH       = Inches(2.25)
COL_GAP         = Inches(0.12)
COL_TOP         = Inches(0.68)
COL_BOTTOM      = Inches(5.05)
COL_HEIGHT      = COL_BOTTOM - COL_TOP

HEADER_H   = Inches(0.48)
TIME_H     = Inches(0.22)
LAB_H      = Inches(0.54)
TOPICS_H   = COL_HEIGHT - HEADER_H - TIME_H - LAB_H - Inches(0.04)

CORNER     = 30000  # adj fmla value for rounded corners

# Column definitions
DAYS = [
    {
        'col_idx': 0,
        'color': '#4A90D9',
        'header1': 'Day 1  |  Step 1',
        'header2': 'LLM과 대화하기',
        'time': '7시간  |  이론 55% · 실습 45%',
        'topics': [
            ('AI 앱 개요: ', '3대 필요성 · 5대 핵심 구성요소',
             '  AI 앱 아키텍처와 개발 흐름'),
            ('LLM API: ', 'OpenAI · Claude · Gemini · Groq',
             '  토큰 비용 · Rate Limit 비교'),
            ('프롬프트 엔지니어링: ', 'Zero/Few-shot · CoT',
             '  체계적 프롬프트 설계'),
            ('Python 기초 & 환경설정:', '',
             '  코드 구성요소 · 기본 문법\n  가상환경 구성'),
            ('멀티턴 대화:', '',
             '  CLIENT/SERVER 상태관리\n  슬라이딩 윈도우 · Responses API'),
        ],
        'lab': '핵심 실습:\n여행플래너 챗봇\n(멀티턴 + Streamlit 앱 UI)',
    },
    {
        'col_idx': 1,
        'color': '#45B869',
        'header1': 'Day 2  |  Step 2~3',
        'header2': '데이터·도구 연결',
        'time': '7시간  |  이론 45% · 실습 55%',
        'topics': [
            ('문서 요약: ', '추출적 / 생성적 요약',
             '  LLM 기반 · 전문 모델 이용'),
            ('PDF 처리 & 멀티모달: ', 'PyMuPDF · IBM Docling',
             '  STT · TTS · 화자분리'),
            ('VLM: ', 'Vision Language Model',
             '  Gemini · Qwen 이미지 분석'),
            ('Function Calling:', '',
             '  도구 호출 워크플로우\n  병렬/순차 · 스트리밍'),
            ('빅3 LLM 고유기능:', '',
             '  OpenAI 커스텀툴 · Claude 서버툴\n  Gemini 자동 함수호출'),
        ],
        'lab': '핵심 실습:\n빅3 LLM 여행플래너\n(Function Calling)',
    },
    {
        'col_idx': 2,
        'color': '#F5A623',
        'header1': 'Day 3~4  |  Step 4',
        'header2': '지식 확장하기',
        'time': '14시간 (7h×2일)  |  이론 63% · 실습 37%',
        'topics': [
            ('LangChain 프레임워크: ', '빠른 개발 · 벤더 독립성',
             '  Workflow · LLM · RAG · Tools'),
            ('RAG 기초: ', '인덱싱 · 검색 메커니즘 · 아키텍처',
             '  RAGAS 지표 평가'),
            ('RAG 품질 튜닝: ', '청킹 · 임베딩 · 하이브리드 검색',
             '  리랭킹 · 쿼리 최적화'),
            ('Local LLM: ', 'Ollama · 양자화 · GPU 레이어',
             '  모델 선택 · 성능 최적화'),
            ('GraphRAG: ', 'MS GraphRAG · Neo4j · LightRAG',
             '  멀티홉 추론 · KG 구축'),
            ('웹검색 + YouTube: ', 'Tavily · DuckDuckGo',
             '  멀티소스 Agent 통합'),
        ],
        'lab': '핵심 실습:\n멀티소스 Agent (웹+YouTube)\nGraphRAG 구축',
    },
    {
        'col_idx': 3,
        'color': '#E74C6F',
        'header1': 'Day 5  |  Step 5',
        'header2': '에이전트 만들기',
        'time': '7시간  |  이론 65% · 실습 35%',
        'topics': [
            ('MAS 아키텍처: ', 'SAS 패턴 · 상태관리 기법',
             '  MAS 내/외부 통신 방법'),
            ('LangGraph: ', '구성요소 · Workflow 개발',
             '  운영 안정성 방안'),
            ('MCP: ', 'Host-Client-Server 아키텍처',
             '  Tools · Resources · Prompts'),
            ('자연어 Agent 개발:', '',
             '  Dify 기반 워크플로우 개발\n  요구사항 → DSL → 프로토타이핑'),
            ('DMAP & Abra:', '',
             '  선언형 멀티에이전트 플러그인\n  Dify 연동 AI 앱 개발'),
        ],
        'lab': '핵심 실습:\nMAS 챗봇 · MCP 서버 구현\nDify 프로토타이핑',
    },
]

for day in DAYS:
    cidx = day['col_idx']
    col_left = COL_LEFT_MARGIN + cidx * (COL_WIDTH + COL_GAP)
    day_color = rgb(day['color'])

    # ── Outer border (rounded rect, transparent fill) ─────────────────────────
    border = rounded_rect(
        slide,
        col_left, COL_TOP, COL_WIDTH, COL_HEIGHT,
        fill_color=None,
        line_color=day_color,
        line_pt=1.5,
        corner_radius_emu=CORNER,
    )

    # ── Header (colored rounded rect covering top portion) ────────────────────
    hdr = rounded_rect(
        slide,
        col_left, COL_TOP, COL_WIDTH, HEADER_H,
        fill_color=day_color,
        line_color=None,
        corner_radius_emu=CORNER,
    )
    # Add header text box on top
    hdr_tx, hdr_tf = add_textbox(
        slide,
        col_left + Inches(0.06),
        COL_TOP + Inches(0.04),
        COL_WIDTH - Inches(0.12),
        HEADER_H - Inches(0.08),
    )
    hdr_tf.word_wrap = True
    for _ in range(len(hdr_tf.paragraphs)):
        hdr_tf._txBody.remove(hdr_tf._txBody.findall(qn('a:p'))[-1])

    ph1 = hdr_tf.add_paragraph()
    ph1.alignment = PP_ALIGN.CENTER
    rh1 = ph1.add_run()
    rh1.text = day['header1']
    rh1.font.bold = False
    rh1.font.size = Pt(7)
    rh1.font.color.rgb = rgb('#FFFFFF')
    rh1.font.name = 'Trebuchet MS'

    ph2 = hdr_tf.add_paragraph()
    ph2.alignment = PP_ALIGN.CENTER
    ph2.space_before = Pt(1)
    rh2 = ph2.add_run()
    rh2.text = day['header2']
    rh2.font.bold = True
    rh2.font.size = Pt(11)
    rh2.font.color.rgb = rgb('#FFFFFF')
    rh2.font.name = 'Trebuchet MS'

    # ── Time info ─────────────────────────────────────────────────────────────
    time_top = COL_TOP + HEADER_H
    time_tx, time_tf = add_textbox(
        slide,
        col_left, time_top, COL_WIDTH, TIME_H,
    )
    time_tf.word_wrap = False
    for _ in range(len(time_tf.paragraphs)):
        time_tf._txBody.remove(time_tf._txBody.findall(qn('a:p'))[-1])
    pt = time_tf.add_paragraph()
    pt.alignment = PP_ALIGN.CENTER
    rt = pt.add_run()
    rt.text = day['time']
    rt.font.bold = False
    rt.font.size = Pt(6)
    rt.font.color.rgb = rgb('#AAAAAA')
    rt.font.name = 'Trebuchet MS'

    # ── Topics text box ───────────────────────────────────────────────────────
    topics_top = COL_TOP + HEADER_H + TIME_H
    topics_tx, topics_tf = add_textbox(
        slide,
        col_left + Inches(0.06),
        topics_top + Inches(0.02),
        COL_WIDTH - Inches(0.12),
        TOPICS_H,
    )
    topics_tf.word_wrap = True
    topics_tf.auto_size = None

    # remove default empty para
    for _ in range(len(topics_tf.paragraphs)):
        topics_tf._txBody.remove(topics_tf._txBody.findall(qn('a:p'))[-1])

    first = True
    for (bold_part, normal_part, continuation) in day['topics']:
        # Title line
        p = topics_tf.add_paragraph()
        p.space_before = Pt(0 if first else 4)
        p.alignment = PP_ALIGN.LEFT
        first = False

        if bold_part:
            rb = p.add_run()
            rb.text = bold_part
            rb.font.bold = True
            rb.font.size = Pt(6.5)
            rb.font.color.rgb = rgb('#E8E8E8')
            rb.font.name = 'Trebuchet MS'
        if normal_part:
            rn = p.add_run()
            rn.text = normal_part
            rn.font.bold = False
            rn.font.size = Pt(5.5)
            rn.font.color.rgb = rgb('#AAAAAA')
            rn.font.name = 'Trebuchet MS'

        # continuation lines (may contain \n)
        for cont_line in continuation.split('\n'):
            pc = topics_tf.add_paragraph()
            pc.space_before = Pt(0)
            pc.alignment = PP_ALIGN.LEFT
            rc = pc.add_run()
            rc.text = cont_line
            rc.font.bold = False
            rc.font.size = Pt(5.5)
            rc.font.color.rgb = rgb('#AAAAAA')
            rc.font.name = 'Trebuchet MS'

    # ── Lab section (bottom colored rounded rect) ─────────────────────────────
    lab_top = COL_BOTTOM - LAB_H
    lab_shape = rounded_rect(
        slide,
        col_left, lab_top, COL_WIDTH, LAB_H,
        fill_color=day_color,
        line_color=None,
        corner_radius_emu=CORNER,
    )
    lab_tx, lab_tf = add_textbox(
        slide,
        col_left + Inches(0.06),
        lab_top + Inches(0.04),
        COL_WIDTH - Inches(0.12),
        LAB_H - Inches(0.08),
    )
    lab_tf.word_wrap = True
    for _ in range(len(lab_tf.paragraphs)):
        lab_tf._txBody.remove(lab_tf._txBody.findall(qn('a:p'))[-1])

    lab_lines = day['lab'].split('\n')
    for li, line in enumerate(lab_lines):
        pl = lab_tf.add_paragraph()
        pl.alignment = PP_ALIGN.CENTER
        pl.space_before = Pt(0 if li > 0 else 0)
        rl = pl.add_run()
        rl.text = line
        rl.font.bold = (li == 0)  # first line bold
        rl.font.size = Pt(6.5)
        rl.font.color.rgb = rgb('#FFFFFF')
        rl.font.name = 'Trebuchet MS'

# ── Arrows between columns ────────────────────────────────────────────────────
for i in range(3):
    arrow_left = (
        COL_LEFT_MARGIN
        + (i + 1) * COL_WIDTH
        + i * COL_GAP
        + Inches(0.01)
    )
    arrow_top  = COL_TOP + Inches(0.12)
    arrow_w    = COL_GAP
    arrow_h    = Inches(0.24)

    ax, atf = add_textbox(slide, arrow_left, arrow_top, arrow_w, arrow_h)
    atf.word_wrap = False
    for _ in range(len(atf.paragraphs)):
        atf._txBody.remove(atf._txBody.findall(qn('a:p'))[-1])
    pa = atf.add_paragraph()
    pa.alignment = PP_ALIGN.CENTER
    ra = pa.add_run()
    ra.text = '▶'
    ra.font.size = Pt(12)
    ra.font.color.rgb = rgb('#666666')
    ra.font.name = 'Trebuchet MS'

# ── Footer ────────────────────────────────────────────────────────────────────
footer_top = Inches(5.30)
footer_h   = Inches(0.28)
footer_tx, footer_tf = add_textbox(
    slide, Inches(0.15), footer_top, Inches(9.70), footer_h
)
footer_tf.word_wrap = False
for _ in range(len(footer_tf.paragraphs)):
    footer_tf._txBody.remove(footer_tf._txBody.findall(qn('a:p'))[-1])

pf = footer_tf.add_paragraph()
pf.alignment = PP_ALIGN.CENTER
rf = pf.add_run()
rf.text = (
    "교재: develop-agent (17개 챕터)  |  "
    "실습 예제: examples/ 디렉토리  |  "
    "전 과정 실습 비중 42%  |  "
    "교육 기간: 5일 (총 35시간)"
)
rf.font.bold = False
rf.font.size = Pt(6.5)
rf.font.color.rgb = rgb('#888888')
rf.font.name = 'Trebuchet MS'

# ── Save ──────────────────────────────────────────────────────────────────────
out_path = r"C:/Users/hiond/workspace/dmap/curriculum_slide.pptx"
prs.save(out_path)
print(f"\nSaved to: {out_path}")
print(f"Total slides: {len(prs.slides)}")
print("Done.")
