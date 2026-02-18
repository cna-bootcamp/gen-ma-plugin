"""
PPTX 슬라이드 추가 스크립트 (재작성)
5일 커리큘럼 연관관계도 슬라이드를 기존 PPTX에 추가
- 흰색(라이트) 테마
- 4컬럼 구조: Day1/Day2/Day3~4/Day5
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ─── 색상 상수 ─────────────────────────────────────────────────
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK_TITLE  = RGBColor(0x1A, 0x1A, 0x1A)
TEAL         = RGBColor(0x22, 0xD3, 0xEE)
GRAY_ARROW   = RGBColor(0x95, 0xA5, 0xA6)
GRAY_MUTED   = RGBColor(0x7F, 0x8C, 0x8D)
DARK_TEXT    = RGBColor(0x1E, 0x29, 0x3B)
DESC_TEXT    = RGBColor(0x47, 0x55, 0x69)

# 컬럼 헤더 색상
C1 = RGBColor(0x4A, 0x90, 0xD9)  # Day 1 파랑
C2 = RGBColor(0x45, 0xB8, 0x69)  # Day 2 초록
C3 = RGBColor(0xF5, 0xA6, 0x23)  # Day 3~4 주황
C4 = RGBColor(0xE7, 0x4C, 0x6F)  # Day 5 빨강

COL_COLORS = [C1, C2, C3, C4]

# ─── 경로 ─────────────────────────────────────────────────────
PPTX_PATH = "C:/Users/hiond/OneDrive/문서/제안/AI Agent 개발/YET-AI Agent-Bootcamp-proposal.pptx"

# ─── 슬라이드 치수 ────────────────────────────────────────────
SLIDE_W = Inches(10.00)
SLIDE_H = Inches(5.62)


# ─── 헬퍼: 마지막 슬라이드 삭제 ──────────────────────────────
def delete_last_slide(prs):
    xml_slides = prs.slides._sldIdLst
    last = xml_slides[-1]
    xml_slides.remove(last)


# ─── 헬퍼: 선 없애기 ──────────────────────────────────────────
def no_line(shape):
    shape.line.fill.background()


# ─── 헬퍼: 사각형 추가 ────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_pt=0.75):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, x, y, w, h)  # 1 = RECTANGLE
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        no_line(shape)
    shape.shadow.inherit = False
    return shape


# ─── 헬퍼: 둥근 사각형 추가 ──────────────────────────────────
def add_rrect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_pt=0.75, adj=20000):
    shape = slide.shapes.add_shape(5, x, y, w, h)  # 5 = ROUNDED_RECTANGLE
    # 둥글기 조정
    try:
        prstGeom = shape._element.spPr.prstGeom
        if prstGeom is not None:
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            for gd in avLst.findall(qn('a:gd')):
                avLst.remove(gd)
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {adj}')
    except Exception:
        pass

    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        no_line(shape)
    shape.shadow.inherit = False
    return shape


# ─── 헬퍼: 텍스트박스 추가 ───────────────────────────────────
def add_tb(slide, x, y, w, h, text, fname="Trebuchet MS", fsize=Pt(8),
           bold=False, rgb=BLACK_TITLE, align=PP_ALIGN.LEFT, wrap=True,
           margin_lr=Pt(2), margin_tb=Pt(1)):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = margin_lr
    tf.margin_right = margin_lr
    tf.margin_top = margin_tb
    tf.margin_bottom = margin_tb
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = fname
    run.font.size = fsize
    run.font.bold = bold
    run.font.color.rgb = rgb
    return txb


# ─── 헬퍼: 줄간격 설정 ───────────────────────────────────────
def set_line_spacing(para, pct=90):
    pPr = para._p.get_or_add_pPr()
    # 기존 lnSpc 제거
    for old in pPr.findall(qn('a:lnSpc')):
        pPr.remove(old)
    lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
    spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
    spcPct.set('val', str(int(pct * 1000)))


# ─── 헬퍼: 단락 전 간격 설정 ─────────────────────────────────
def set_space_before(para, pt_val):
    pPr = para._p.get_or_add_pPr()
    for old in pPr.findall(qn('a:spcBef')):
        pPr.remove(old)
    spcBef = etree.SubElement(pPr, qn('a:spcBef'))
    spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
    spcPts.set('val', str(int(pt_val * 100)))


# ─── 컬럼 데이터 ──────────────────────────────────────────────
COLUMNS = [
    {
        "color": C1,
        "header_line1": "Day 1 | Step 1",
        "header_line2": "LLM과 대화하기",
        "time": "7시간  |  이론 55% · 실습 45%",
        "topics": [
            ("AI 앱 개요", "3대 필요성 · 5대 핵심 구성요소\nAI 앱 아키텍처와 개발 흐름"),
            ("LLM API", "OpenAI · Claude · Gemini · Groq\n토큰 비용 · Rate Limit 비교"),
            ("프롬프트 엔지니어링", "Zero/Few-shot · CoT\n체계적 프롬프트 설계"),
            ("Python 기초 & 환경설정", "코드 구성요소 · 기본 문법\n가상환경 구성"),
            ("멀티턴 대화", "CLIENT/SERVER 상태관리\n슬라이딩 윈도우 · Responses API"),
        ],
        "lab": "핵심 실습:\n여행플래너 챗봇\n(멀티턴 + Streamlit 앱 UI)",
    },
    {
        "color": C2,
        "header_line1": "Day 2 | Step 2~3",
        "header_line2": "데이터·도구 연결",
        "time": "7시간  |  이론 45% · 실습 55%",
        "topics": [
            ("문서 요약", "추출적 / 생성적 요약\nLLM 기반 · 전문 모델 이용"),
            ("PDF 처리 & 멀티모달", "PyMuPDF · IBM Docling\nSTT · TTS · 화자분리"),
            ("VLM", "Vision Language Model\nGemini · Qwen 이미지 분석"),
            ("Function Calling", "도구 호출 워크플로우\n병렬/순차 · 스트리밍"),
            ("빅3 LLM 고유기능", "OpenAI 커스텀툴 · Claude 서버툴\nGemini 자동 함수호출"),
        ],
        "lab": "핵심 실습:\n빅3 LLM 여행플래너\n(Function Calling)",
    },
    {
        "color": C3,
        "header_line1": "Day 3~4 | Step 4",
        "header_line2": "지식 확장하기",
        "time": "14시간 (7h×2일)  |  이론 63% · 실습 37%",
        "topics": [
            ("LangChain 프레임워크", "빠른 개발 · 벤더 독립성\nWorkflow · LLM · RAG · Tools"),
            ("RAG 기초", "인덱싱 · 검색 메커니즘 · 아키텍처\nRAGAS 지표 평가"),
            ("RAG 품질 튜닝", "청킹 · 임베딩 · 하이브리드 검색\n리랭킹 · 쿼리 최적화"),
            ("Local LLM", "Ollama · 양자화 · GPU 레이어\n모델 선택 · 성능 최적화"),
            ("GraphRAG", "MS GraphRAG · Neo4j · LightRAG\n멀티홉 추론 · KG 구축"),
            ("웹검색 + YouTube", "Tavily · DuckDuckGo\n멀티소스 Agent 통합"),
        ],
        "lab": "핵심 실습:\n멀티소스 Agent (웹+YouTube)\nGraphRAG 구축",
    },
    {
        "color": C4,
        "header_line1": "Day 5 | Step 5",
        "header_line2": "에이전트 만들기",
        "time": "7시간  |  이론 65% · 실습 35%",
        "topics": [
            ("MAS 아키텍처", "SAS 패턴 · 상태관리 기법\nMAS 내/외부 통신 방법"),
            ("LangGraph", "구성요소 · Workflow 개발\n운영 안정성 방안"),
            ("MCP", "Host-Client-Server 아키텍처\nTools · Resources · Prompts"),
            ("자연어 Agent 개발", "Dify 기반 워크플로우 개발\n요구사항 → DSL → 프로토타이핑"),
            ("DMAP & Abra", "선언형 멀티에이전트 플러그인\nDify 연동 AI 앱 개발"),
        ],
        "lab": "핵심 실습:\nMAS 챗봇 · MCP 서버 구현\nDify 프로토타이핑",
    },
]

# ─── 메인 ─────────────────────────────────────────────────────

prs = Presentation(PPTX_PATH)
print(f"기존 슬라이드 수: {len(prs.slides)}")

# 슬라이드가 15개이면 마지막 슬라이드 삭제
if len(prs.slides) == 15:
    delete_last_slide(prs)
    print(f"마지막 슬라이드 삭제 완료 → 현재 슬라이드 수: {len(prs.slides)}")

# 빈 레이아웃(blank)으로 슬라이드 추가
# 레이아웃 6번이 보통 blank, 없으면 0번 사용
blank_layout = None
for layout in prs.slide_layouts:
    if layout.name.lower() in ("blank", "blank slide", "빈 화면", ""):
        blank_layout = layout
        break
if blank_layout is None:
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

slide = prs.slides.add_slide(blank_layout)

# placeholder 모두 제거
for ph in list(slide.placeholders):
    ph._element.getparent().remove(ph._element)

W = prs.slide_width   # Inches(10)
H = prs.slide_height  # Inches(5.62)

# ─── 1. 흰색 배경 ─────────────────────────────────────────────
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

# ─── 2. 제목 ──────────────────────────────────────────────────
title_y = Inches(0.12)
add_tb(slide, Inches(0.3), title_y, W - Inches(0.6), Inches(0.30),
       "YET AI Agent 부트캠프 — 5일 커리큘럼 연관관계도",
       fname="Arial Black", fsize=Pt(17), bold=True, rgb=BLACK_TITLE,
       align=PP_ALIGN.CENTER, margin_lr=Pt(4), margin_tb=Pt(2))

subtitle_y = Inches(0.46)
add_tb(slide, Inches(0.3), subtitle_y, W - Inches(0.6), Inches(0.22),
       "대화  ▶  데이터  ▶  도구  ▶  지식  ▶  에이전트",
       fname="Trebuchet MS", fsize=Pt(10), bold=False, rgb=TEAL,
       align=PP_ALIGN.CENTER, margin_lr=Pt(4), margin_tb=Pt(1))

# ─── 3. 4개 컬럼 레이아웃 계산 ───────────────────────────────
# 전체 높이 배분:
#   타이틀 영역: 0~0.70 in
#   컬럼 시작: 0.72 in
#   컬럼 끝 (footer 위): H - 0.28 in
#   footer: 0.25 in

COL_TOP = Inches(0.72)
FOOTER_H = Inches(0.24)
COL_BOT = H - FOOTER_H - Inches(0.02)
COL_TOTAL_H = COL_BOT - COL_TOP

# 헤더 높이
HDR_H = Inches(0.48)
# 시간정보 높이
TIME_H = Inches(0.18)
# 실습 섹션 높이
LAB_H = Inches(0.44)
# 토픽 영역 높이 (나머지)
TOPIC_H = COL_TOTAL_H - HDR_H - TIME_H - LAB_H - Inches(0.04)

# 컬럼 너비
LEFT_MARGIN = Inches(0.25)
RIGHT_MARGIN = Inches(0.25)
ARROW_W = Inches(0.18)
NUM_GAPS = 3
available = W - LEFT_MARGIN - RIGHT_MARGIN - ARROW_W * NUM_GAPS
COL_W = available / 4

# 컬럼 x 위치 계산 (화살표 공간 포함)
col_xs = []
x = LEFT_MARGIN
for i in range(4):
    col_xs.append(x)
    x += COL_W
    if i < 3:
        x += ARROW_W

# ─── 화살표 (컬럼 사이) ──────────────────────────────────────
arrow_y = COL_TOP + HDR_H / 2 - Inches(0.10)
for i in range(3):
    ax = col_xs[i] + COL_W
    add_tb(slide, ax, arrow_y, ARROW_W, Inches(0.20),
           "▶", fname="Trebuchet MS", fsize=Pt(10), bold=False,
           rgb=GRAY_ARROW, align=PP_ALIGN.CENTER,
           margin_lr=Pt(0), margin_tb=Pt(0))

# ─── 컬럼 그리기 ──────────────────────────────────────────────
for ci, col in enumerate(COLUMNS):
    cx = col_xs[ci]
    cy = COL_TOP
    color = col["color"]

    # 컬럼 전체 테두리 (둥근 사각형, 배경 없음, 컬러 테두리)
    border = add_rrect(slide, cx, cy, COL_W, COL_TOTAL_H,
                       fill_rgb=None, line_rgb=color, line_pt=1.0, adj=15000)

    # ── 헤더 배경 (둥근 사각형, 컬럼 색상) ──────────────────
    hdr_shape = add_rrect(slide, cx, cy, COL_W, HDR_H,
                          fill_rgb=color, line_rgb=None, adj=15000)

    # 헤더 텍스트 (shape 텍스트 프레임 사용)
    tf = hdr_shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(5)
    tf.margin_bottom = Pt(2)

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = col["header_line1"]
    r1.font.name = "Trebuchet MS"
    r1.font.size = Pt(7)
    r1.font.bold = False
    r1.font.color.rgb = WHITE
    set_line_spacing(p1, 90)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = col["header_line2"]
    r2.font.name = "Trebuchet MS"
    r2.font.size = Pt(10)
    r2.font.bold = True
    r2.font.color.rgb = WHITE
    set_line_spacing(p2, 90)

    # ── 시간 정보 ─────────────────────────────────────────────
    time_y = cy + HDR_H
    add_tb(slide, cx + Pt(2), time_y, COL_W - Pt(4), TIME_H,
           col["time"],
           fname="Trebuchet MS", fsize=Pt(6), bold=False, rgb=GRAY_MUTED,
           align=PP_ALIGN.CENTER, margin_lr=Pt(2), margin_tb=Pt(2))

    # ── 토픽 영역 텍스트박스 ──────────────────────────────────
    topic_y = time_y + TIME_H
    pad_lr = Pt(5)
    topic_txb = slide.shapes.add_textbox(
        cx + pad_lr, topic_y,
        COL_W - pad_lr * 2, TOPIC_H
    )
    tf2 = topic_txb.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Pt(2)
    tf2.margin_right = Pt(2)
    tf2.margin_top = Pt(3)
    tf2.margin_bottom = Pt(2)

    first_para = True
    for (title, desc) in col["topics"]:
        if first_para:
            para = tf2.paragraphs[0]
            first_para = False
        else:
            para = tf2.add_paragraph()
            set_space_before(para, 3.0)

        para.alignment = PP_ALIGN.LEFT
        set_line_spacing(para, 90)

        # 제목 run (bold, 어두운 색)
        r_title = para.add_run()
        r_title.text = title + ": "
        r_title.font.name = "Trebuchet MS"
        r_title.font.size = Pt(6.5)
        r_title.font.bold = True
        r_title.font.color.rgb = DARK_TEXT

        # 설명 첫 줄이 있으면 같은 단락에 추가
        desc_lines = desc.split("\n")
        if desc_lines and desc_lines[0].strip():
            r_desc1 = para.add_run()
            r_desc1.text = desc_lines[0]
            r_desc1.font.name = "Trebuchet MS"
            r_desc1.font.size = Pt(6)
            r_desc1.font.bold = False
            r_desc1.font.color.rgb = DESC_TEXT

        # 설명 두 번째 줄 이후
        for dl in desc_lines[1:]:
            if dl.strip():
                p_sub = tf2.add_paragraph()
                p_sub.alignment = PP_ALIGN.LEFT
                set_line_spacing(p_sub, 88)
                # 들여쓰기 효과를 위해 공백 추가
                r_sub = p_sub.add_run()
                r_sub.text = "  " + dl
                r_sub.font.name = "Trebuchet MS"
                r_sub.font.size = Pt(5.8)
                r_sub.font.bold = False
                r_sub.font.color.rgb = DESC_TEXT

    # ── 실습 섹션 (하단) ──────────────────────────────────────
    lab_y = cy + COL_TOTAL_H - LAB_H
    lab_shape = add_rect(slide, cx, lab_y, COL_W, LAB_H,
                         fill_rgb=color, line_rgb=None)

    tf3 = lab_shape.text_frame
    tf3.word_wrap = True
    tf3.margin_left = Pt(4)
    tf3.margin_right = Pt(4)
    tf3.margin_top = Pt(4)
    tf3.margin_bottom = Pt(2)

    lab_lines = col["lab"].split("\n")
    for li, lab_line in enumerate(lab_lines):
        if li == 0:
            lp = tf3.paragraphs[0]
        else:
            lp = tf3.add_paragraph()
        lp.alignment = PP_ALIGN.CENTER
        set_line_spacing(lp, 88)
        lr = lp.add_run()
        lr.text = lab_line
        lr.font.name = "Trebuchet MS"
        lr.font.size = Pt(6.5) if li == 0 else Pt(6)
        lr.font.bold = (li == 0)
        lr.font.color.rgb = WHITE

# ─── 4. 푸터 ──────────────────────────────────────────────────
footer_y = H - FOOTER_H
add_tb(slide, Inches(0.25), footer_y, W - Inches(0.5), FOOTER_H,
       "교재: develop-agent (17개 챕터)  |  실습 예제: examples/ 디렉토리  |  전 과정 실습 비중 42%  |  교육 기간: 5일 (총 35시간)",
       fname="Trebuchet MS", fsize=Pt(7), bold=False, rgb=GRAY_MUTED,
       align=PP_ALIGN.CENTER, margin_lr=Pt(4), margin_tb=Pt(4))

# ─── 저장 ─────────────────────────────────────────────────────
prs.save(PPTX_PATH)
print(f"저장 완료: {PPTX_PATH}")
print(f"최종 슬라이드 수: {len(prs.slides)}")
print("슬라이드 추가 성공: 5일 커리큘럼 연관관계도 (흰색 배경)")
